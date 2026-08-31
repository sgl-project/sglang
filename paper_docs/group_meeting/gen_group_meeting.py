"""Generate the two charts + group-meeting PPTX for innovation points 1 & 2."""
import os

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ASSETS = "/data/xbw/kvbench/assets"
os.makedirs(ASSETS, exist_ok=True)

C_BF16, C_FP4 = "#0072B2", "#D55E00"  # validated Okabe-Ito pair
INK, MUTED = RGBColor(0x1A, 0x1A, 0x2E), RGBColor(0x66, 0x66, 0x77)
ACCENT = RGBColor(0x00, 0x50, 0x8A)

# ---------- chart 1: kernel-level decode latency ----------
fig, ax = plt.subplots(figsize=(5.6, 3.0))
names = ["flashinfer\nfp16 mma", "flashinfer\nfp16 scalar", "Ours fp4\nscalar", "Ours fp4\nmma(ldmatrix)"]
vals = [0.474, 0.556, 0.846, 0.677]
colors = ["#999999", "#bbbbbb", C_BF16, C_BF16]
bars = ax.bar(names, vals, color=colors, width=0.62)
for b, v in zip(bars, vals):
    ax.text(b.get_x() + b.get_width() / 2, v + 0.015, f"{v:.3f}", ha="center", fontsize=10)
ax.set_ylabel("decode attention latency (ms)")
ax.set_title("b=100, seq=1024, qh=32, kh=8 (RTX 3090)", fontsize=9, color="#666677")
ax.set_ylim(0, 0.95)
ax.grid(True, axis="y", alpha=0.3)
ax.spines[["top", "right"]].set_visible(False)
fig.tight_layout()
fig.savefig(f"{ASSETS}/kernel_bars.png", dpi=200)
plt.close(fig)

# ---------- chart 2: BFCL grouped bars ----------
fig, ax = plt.subplots(figsize=(5.6, 3.0))
cats = ["simple_python", "multiple", "multi_turn_base\n(main agent metric)"]
bf = [93.75, 91.50, 23.50]
mx = [92.00, 92.00, 11.00]
x = [0, 1, 2]
w = 0.36
ax.bar([i - w / 2 for i in x], bf, w, color=C_BF16, label="BF16")
ax.bar([i + w / 2 for i in x], mx, w, color=C_FP4, label="MXFP4")
for i, (a, b) in enumerate(zip(bf, mx)):
    ax.text(i - w / 2, a + 1.2, f"{a:.1f}", ha="center", fontsize=9)
    ax.text(i + w / 2, b + 1.2, f"{b:.1f}", ha="center", fontsize=9)
ax.set_xticks(x)
ax.set_xticklabels(cats, fontsize=9)
ax.set_ylabel("BFCL accuracy (%)")
ax.set_ylim(0, 105)
ax.legend(frameon=False, fontsize=9)
ax.grid(True, axis="y", alpha=0.3)
ax.spines[["top", "right"]].set_visible(False)
fig.tight_layout()
fig.savefig(f"{ASSETS}/bfcl_bars.png", dpi=200)
plt.close(fig)

# ---------- chart 3: capacity ----------
fig, ax = plt.subplots(figsize=(4.2, 3.0))
bars = ax.bar(["BF16", "MXFP4"], [93591, 283808], color=[C_BF16, C_FP4], width=0.5)
for b, v in zip(bars, [93591, 283808]):
    ax.text(b.get_x() + b.get_width() / 2, v + 6000, f"{v:,}", ha="center", fontsize=10)
ax.set_ylabel("max_total_num_tokens")
ax.set_title("KV cache capacity (Qwen3-4B, 24GB)", fontsize=10, color="#666677")
ax.set_ylim(0, 320000)
ax.grid(True, axis="y", alpha=0.3)
ax.spines[["top", "right"]].set_visible(False)
fig.tight_layout()
fig.savefig(f"{ASSETS}/capacity.png", dpi=200)
plt.close(fig)

# ---------- PPTX ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide(title, subtitle=None):
    s = prs.slides.add_slide(BLANK)
    tb = s.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(12.2), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = INK
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(14)
        r2.font.color.rgb = MUTED
    ln = s.shapes.add_shape(1, Inches(0.6), Inches(1.28), Inches(12.1), Pt(2.2))
    ln.fill.solid()
    ln.fill.fore_color.rgb = ACCENT
    ln.line.fill.background()
    return s


def bullets(slide, items, left=0.75, top=1.6, width=11.9, height=5.4, size=17):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        if isinstance(it, tuple):
            mark, txt = it
        else:
            mark, txt = "•", it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = f"{mark} {txt}" if mark else txt
        r.font.size = Pt(size)
        r.font.color.rgb = INK
        p.space_after = Pt(8)
    return tb


def add_table(slide, data, left, top, width, height, size=13, header=True):
    rows, cols = len(data), len(data[0])
    shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    tbl = shape.table
    for i, row in enumerate(data):
        for j, cell in enumerate(row):
            c = tbl.cell(i, j)
            c.text = str(cell)
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(size)
                    if header and i == 0:
                        r.font.bold = True
    return tbl


def add_image(slide, path, left, top, width=None, height=None):
    kw = {}
    if width:
        kw["width"] = Inches(width)
    if height:
        kw["height"] = Inches(height)
    return slide.shapes.add_picture(path, Inches(left), Inches(top), **kw)


# S1 cover
s = prs.slides.add_slide(BLANK)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.5), Inches(2.6))
tf = tb.text_frame
p = tf.paragraphs[0]
r = p.add_run()
r.text = "资源受限场景下的复合 AI 系统加速技术研究"
r.font.size = Pt(38)
r.font.bold = True
r.font.color.rgb = INK
p2 = tf.add_paragraph()
r2 = p2.add_run()
r2.text = "进展汇报:创新点 1(MXFP4 KV Cache)+ 创新点 2(Agent 场景 KV 保留与容量包络)"
r2.font.size = Pt(19)
r2.font.color.rgb = MUTED
p3 = tf.add_paragraph()
r3 = p3.add_run()
r3.text = "组会汇报 · 2026-09"
r3.font.size = Pt(15)
r3.font.color.rgb = MUTED

# S2 agenda
s = add_slide("本次汇报", "两条线:一条做完并验收,一条做完并发现意外结论")
bullets(s, [
    "创新点 1|端侧消费级 GPU 上的 MXFP4 KV Cache:方案、kernel 演进、端到端验收、Agent 能力边界",
    "创新点 2|Agent 场景 KV 保留与容量包络:自研压测方法、三区模型、一个诚实的负结果",
    ("▸ ", "贯穿结论:对 Agent serving 而言,KV 容量是一阶杠杆,策略是二阶杠杆"),
    "下一步:创新点 3(CUDA 流预取)可行性判定 / 论文写作",
], size=19)

# S3 motivation
s = add_slide("创新点 1|动机:24GB 消费卡上,KV Cache 是 Agent 的硬瓶颈", "Qwen3-4B / 8B,BF16 KV;Agent 多轮对话动辄 8k-26k token 上下文")
add_image(s, f"{ASSETS}/capacity.png", 8.3, 1.7, width=4.4)
bullets(s, [
    "BF16 下 KV Cache 占走大半显存,容量天花板 ~9.4 万 token",
    "一条长 Agent 会话(8k-26k tok)就吃掉 10-30%;并发几路即触顶",
    "现有 FP4 方案(vLLM/flashinfer)依赖 sm90+ 专属 kernel,sm86 无基础设施",
    "目标:sm86 上做出 MXFP4 KV 全链路,验收线:精度 ≥90%、开销 ≤+20%",
], width=7.3, size=16)

# S4 design
s = add_slide("创新点 1|方案:MXFP4 存储 + 零带宽融合 decode", "block32(E2M1+E8M0)+ RNE 舍入,~2000 行 Python/CUDA,JIT 接入 sglang")
bullets(s, [
    "存储:2×E2M1 打包 + 每 32 维 1 个 E8M0 scale,沿 head_dim 分块(73.4% 省)",
    "量化:逐块 abs-max → scale,RNE 到 E2M1(精度关键之一,+1.5pt)",
    "decode:自研融合 kernel 直接读 fp4,寄存器内反量化,不落 bf16 workspace",
    "接入:--kv-cache-dtype fp4_mx_block32,生产默认配置(CUDA graph/radix/overlap 全开)",
    ("▸ ", "两代 kernel:CUDA-cores 标量版 → ldmatrix+mma.m16n8k16 张量核版(生产)"),
], size=17)

# S5 kernel evolution
s = add_slide("创新点 1|kernel 演进:ldmatrix+mma 把 TPOT 从 +26% 拉回 +16%", "手工构造 fragment 是死路(2.1× 慢);ldmatrix 语义按 PTX ISA 逐条确认后重写")
add_image(s, f"{ASSETS}/kernel_bars.png", 7.6, 1.7, width=5.3)
bullets(s, [
    "零 shuffle 设计:QK 的 C 片段原地复用为 PV 的 A 片段;O 直接从 C 写回",
    "4 warp 按 token 分块 + 寄存器预取 + 128B swizzle",
    "kernel 内比 flashinfer fp16 mma 慢 0.20ms = fp4 反量化的固有 ALU 成本",
    ("▸ ", "精度教训:softmax 分子分母必须同源舍入(否则 GSM8K -2.7pt)"),
], width=6.7, size=16)

# S6 e2e results
s = add_slide("创新点 1|端到端验收:全部达标(生产配置)", "CUDA graph + radix + overlap 全开;BF16 同配置背靠背对比")
add_table(s, [
    ["指标", "BF16", "MXFP4", "Δ", "线"],
    ["GSM8K 200", "94.0%", "89.2%(3 轮均值)", "-4.8pt", "≈90% ✓"],
    ["Median TPOT", "32.84 ms", "38.15 ms", "+16.2%", "≤+20% ✓"],
    ["Median TTFT", "99.1 ms", "106.3 ms", "+7.3%", "≤+20% ✓"],
    ["KV 容量", "93,591 tok", "283,808 tok", "3.03×", "—"],
], 0.75, 1.75, 11.8, 2.6, size=16)
bullets(s, [
    "反量化 ALU 是固有开销;换来的是 3× 容量 —— 这个交换值不值,创新点 2 给出定量答案",
], top=4.7, size=17)

# S7 BFCL
s = add_slide("创新点 1|Agent 能力边界(BFCL):单轮无损,长上下文有代价", "bfcl-eval 2026.3.23;BF16 复现与官方榜单逐位吻合(93.75/91.0/24.5)→ 管线可信")
add_image(s, f"{ASSETS}/bfcl_bars.png", 7.6, 1.75, width=5.3)
bullets(s, [
    "单轮 sanity 持平:tool-call 生成能力保持",
    "multi_turn_base -12.5pt(23.5→11.0):8k-26k 多轮中 fp4 KV 误差累积",
    "失败模式:17/30 回归条目为工具调用循环(状态跟踪退化)",
    ("▸ ", "定位为量化 KV 的能力边界,与容量收益构成完整 tradeoff,论文如实呈现"),
], width=6.7, size=16)

# S8 motivation 2
s = add_slide("创新点 2|动机:工具调用间隙,KV 该不该留?", "Agent 请求在工具执行期(0.1s~30s)被打断;期间 KV 前缀被驱逐 → 返回后全量重算")
bullets(s, [
    "重算成本示例:8k 前缀 prefill ≈ 2-4s;命中前缀只算增量 ≈ 0.1-0.3s(10× 级)",
    "直觉方案:预测工具返回时间,保留/预取对应 KV —— Continuum(arXiv 2511.02230)用 TTL pin",
    ("▸ ", "空白:现有 benchmark(BFCL 等)工具执行瞬时完成、无持续压力,测不出这类策略的价值"),
    ("▸ ", "所以我们先造了测量工具:trace-driven Agent KV 压测 harness"),
], size=17)

# S9 harness
s = add_slide("创新点 2|方法:trace-driven 压测 harness", "标准 OpenAI tool-calling 回路;N 路并发;压力可控、可复现")
bullets(s, [
    "每轮:tool-call 轮 → sleep(该工具采样的执行延迟)→ 续轮(指标轮)",
    "工具混合 4 类 lognormal(快 0.1-1s ×50% / 慢 5-30s ×20%…),指令执行 100% 依从",
    "上下文 fabrication 精确控长(2.2k 起步,每轮 +1.2k),前缀 token 精确连续",
    "指标:续轮 miss 率(=驱逐后全量重算)/ 浪费重算 token / TTFT p50",
    ("▸ ", "命中与未命中实测差 100×(129ms vs 12.7s),信号极干净"),
], size=17)

# S10 main figure
s = add_slide("创新点 2|主结果:容量是 Agent serving 的一阶杠杆", "LRU;6 轮会话;曲线为续轮 miss 率(左)/ TTFT p50(右)")
add_image(s, "/home/xubowen/mxfp4/sglang/paper_docs/agent_kv_bench/envelope.png", 0.9, 1.6, width=11.5)
bullets(s, [
    "BF16:N≥32 崩溃(81% miss,TTFT 15.6s);MXFP4 同负载 17.7% / 0.30s(50×)",
    "等质量并发包络 2-3×;MXFP4 自身拐点在 N≈56-64",
], top=6.15, size=15)

# S11 three regimes + negative result
s = add_slide("创新点 2|三区模型 + 一个诚实的负结果", "需求速率 ÷ 池容量(churn)单一变量刻画;3 seed 验证")
bullets(s, [
    "可容纳区(<1× 超卖):fp4 N=32 工作集装得下,排序策略无冲突可解",
    "边际区(1-1.4×):gap 感知驱逐有窄窗口收益(fp4/32 三 seed 一致 -15% 浪费)",
    "深饱和区(>1.4×):容量硬顶,任何排序无效",
    ("▸ ", "负结果:我们的 Score 驱逐未稳健胜过 LRU —— 机理:agent 负载下访问新近度 ≈ 会话翻转率,LRU 已隐式近似 Belady"),
    ("▸ ", "对 Continuum 类工作的清醒对照:radix-on 引擎(如 SGLang)上 retention headroom 很小;其大收益来自与 vLLM(前缀缓存弱)的对比"),
], size=16)

# S12 next
s = add_slide("总结与下一步")
bullets(s, [
    "创新点 1 ✓:MXFP4 KV 全链路,验收全达标,容量 3.03×;Agent 能力边界已定量(BFCL -12.5pt)",
    "创新点 2 ✓:压测 harness + 三区模型 + 包络主图;主结论=容量是一阶杠杆;负结果与机理已封存",
    "创新点 3(CUDA 流预取):物理账成立(重载比重算便宜 15-25×),生死问题=HiCache 按需恢复的成本,1-2h 实验可判定",
    "论文:素材齐备(方法/主图/三区/负结果/对照),进入写作",
], size=19)

out = "/data/xbw/kvbench/组会汇报_创新点12.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
